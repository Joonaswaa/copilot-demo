"""
pptx_exporter.py
----------------
Builds a weekly management deck from dashboard analytics,
optionally enriched with narrative sections from the weekly report (hybrid mode).
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from src.report_parsing import parse_executive_summary, parse_management_actions

# Brand colours (aligned with example decks + dashboard)
NAVY = RGBColor(0x1E, 0x27, 0x61)
ACCENT = RGBColor(0x00, 0x2C, 0xAC)
KPI_BG = RGBColor(0xF1, 0xF5, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1D, 0x1D, 0x1F)
MUTED = RGBColor(0x6E, 0x6E, 0x73)
RISK_COLORS = {
    "Low": RGBColor(0x2E, 0x9E, 0x5B),
    "Medium": RGBColor(0xE3, 0xA0, 0x08),
    "High": RGBColor(0xE8, 0x59, 0x0C),
    "Critical": RGBColor(0xC9, 0x2A, 0x2A),
}
RISK_LEVELS = ["Low", "Medium", "High", "Critical"]

LABELS = {
    "en": {
        "title": "Weekly Supply Chain Review",
        "subtitle": "Telecom Supply Chain AI Copilot — management summary",
        "kpi_title": "This week in numbers",
        "kpi_inventory": "Inventory value",
        "kpi_skus": "SKUs",
        "kpi_high_risk": "High-risk products",
        "kpi_delayed": "Delayed deliveries",
        "kpi_purchases": "Recommended purchases",
        "kpi_on_time": "Supplier on-time rate",
        "risk_title": "Risk overview",
        "risk_top": "Highest-risk products",
        "weeks_cover": "{weeks:.1f} weeks of cover · {supplier}",
        "purchase_title": "Recommended purchase orders",
        "col_product": "Product",
        "col_supplier": "Supplier",
        "col_qty": "Qty",
        "col_cost": "Cost (€)",
        "col_priority": "Priority",
        "purchase_total": "Total proposed purchase value: {value}",
        "supplier_title": "Supplier performance",
        "supplier_attention": "Attention needed",
        "supplier_warn": (
            "•  {supplier}: on-time rate is only {on_time:.0f}%, "
            "average delay is {delay:.1f} days, {delayed} deliveries are "
            "currently delayed. Affected categories: {categories}."
        ),
        "supplier_warn_no_delay": (
            "•  {supplier}: on-time rate is only {on_time:.0f}%, "
            "{delayed} deliveries are currently delayed. "
            "Affected categories: {categories}."
        ),
        "supplier_warn_delay_only": (
            "•  {supplier}: on-time rate is only {on_time:.0f}%, "
            "average delay is {delay:.1f} days. "
            "Affected categories: {categories}."
        ),
        "slow_title": "Slow-moving inventory",
        "slow_capital": "capital tied in {count} slow-moving SKUs",
        "actions_title": "Recommended actions",
        "exec_summary_title": "Executive summary",
        "exec_summary_note": "Narrative from the weekly AI report — figures match the KPI slide.",
    },
    "fi": {
        "title": "Viikoittainen toimitusketjukatsaus",
        "subtitle": "Telecom Supply Chain AI Copilot — johdon yhteenveto",
        "kpi_title": "Viikko numeroina",
        "kpi_inventory": "Varaston arvo",
        "kpi_skus": "Nimikkeet",
        "kpi_high_risk": "Korkean riskin tuotteet",
        "kpi_delayed": "Myöhässä olevat toimitukset",
        "kpi_purchases": "Suositellut ostot",
        "kpi_on_time": "Toimitusvarmuus",
        "risk_title": "Riskitilanne",
        "risk_top": "Korkeimman riskin tuotteet",
        "weeks_cover": "{weeks:.1f} viikon riitto · {supplier}",
        "purchase_title": "Suositellut ostotilaukset",
        "col_product": "Tuote",
        "col_supplier": "Toimittaja",
        "col_qty": "Määrä",
        "col_cost": "Kustannus (€)",
        "col_priority": "Prioriteetti",
        "purchase_total": "Ehdotettujen ostojen kokonaisarvo: {value}",
        "supplier_title": "Toimittajien suorituskyky",
        "supplier_attention": "Vaatii huomiota",
        "supplier_warn": (
            "•  {supplier}: toimitusvarmuus on vain {on_time:.0f} %, "
            "keskimääräinen viive on {delay:.1f} päivää, {delayed} toimitusta "
            "on tällä hetkellä myöhässä. Kategoriat: {categories}."
        ),
        "supplier_warn_no_delay": (
            "•  {supplier}: toimitusvarmuus on vain {on_time:.0f} %, "
            "{delayed} toimitusta on tällä hetkellä myöhässä. "
            "Kategoriat: {categories}."
        ),
        "supplier_warn_delay_only": (
            "•  {supplier}: toimitusvarmuus on vain {on_time:.0f} %, "
            "keskimääräinen viive on {delay:.1f} päivää. "
            "Kategoriat: {categories}."
        ),
        "slow_title": "Hitaasti kiertävä varasto",
        "slow_capital": "pääomaa sitoutunut {count} hitaasti kiertävään nimikkeeseen",
        "actions_title": "Toimenpidesuositukset",
        "exec_summary_title": "Johdon yhteenveto",
        "exec_summary_note": "Teksti viikoittaisesta AI-raportista — luvut vastaavat KPI-diaa.",
    },
}


def _fmt_eur(value: float) -> str:
    return f"€{value:,.0f}"


def _fmt_date(language: str) -> str:
    today = date.today()
    if language == "fi":
        return today.strftime("%d.%m.%Y")
    return today.strftime("%d.%m.%Y")


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_run(run, *, size: int, bold: bool = False, color: RGBColor = DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _textbox(slide, left, top, width, height, text: str, *,
             size: int = 18, bold: bool = False, color: RGBColor = DARK,
             align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return box


def _fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _slide_title(slide, text: str):
    _textbox(slide, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.8),
             text, size=28, bold=True, color=ACCENT)


def _kpi_card(slide, left, top, value: str, label: str):
    w, h = Inches(3.9), Inches(2.3)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    _fill_shape(card, KPI_BG)
    _textbox(slide, left, top + Inches(0.45), w, Inches(1.0),
             value, size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _textbox(slide, left + Inches(0.2), top + Inches(1.55), w - Inches(0.4),
             Inches(0.6), label, size=13, color=MUTED, align=PP_ALIGN.CENTER)


def _add_title_slide(prs: Presentation, language: str):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _fill_shape(bg, NAVY)
    lbl = LABELS[language]
    _textbox(slide, Inches(0.6), Inches(2.5), Inches(12.1), Inches(1.3),
             lbl["title"], size=40, bold=True, color=WHITE)
    _textbox(slide, Inches(0.6), Inches(3.8), Inches(12.1), Inches(0.6),
             lbl["subtitle"], size=18, color=WHITE)
    _textbox(slide, Inches(0.6), Inches(4.5), Inches(12.1), Inches(0.5),
             _fmt_date(language), size=16, color=WHITE)


def _add_kpi_slide(prs: Presentation, kpis: dict, language: str):
    slide = _blank_slide(prs)
    lbl = LABELS[language]
    _slide_title(slide, lbl["kpi_title"])

    cards = [
        (_fmt_eur(kpis["total_inventory_value"]), lbl["kpi_inventory"]),
        (str(kpis["sku_count"]), lbl["kpi_skus"]),
        (str(kpis["high_risk_count"]), lbl["kpi_high_risk"]),
        (str(kpis["delayed_deliveries"]), lbl["kpi_delayed"]),
        (_fmt_eur(kpis["recommended_purchase_value"]), lbl["kpi_purchases"]),
        (f"{kpis['avg_on_time_rate']:.1f}%", lbl["kpi_on_time"]),
    ]
    xs = [Inches(0.47), Inches(4.72), Inches(8.97)]
    ys = [Inches(1.7), Inches(4.35)]
    for i, (value, label) in enumerate(cards):
        _kpi_card(slide, xs[i % 3], ys[i // 3], value, label)


def _add_risk_slide(prs: Presentation, df: pd.DataFrame, language: str):
    slide = _blank_slide(prs)
    lbl = LABELS[language]
    _slide_title(slide, lbl["risk_title"])

    counts = df["risk_level"].value_counts().reindex(RISK_LEVELS).fillna(0)
    chart_data = CategoryChartData()
    chart_data.categories = RISK_LEVELS
    chart_data.add_series("Risk", counts.tolist())

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(0.6), Inches(1.6), Inches(5.6), Inches(5.2),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    for i, level in enumerate(RISK_LEVELS):
        if i < len(plot.series[0].points):
            plot.series[0].points[i].format.fill.solid()
            plot.series[0].points[i].format.fill.fore_color.rgb = RISK_COLORS[level]

    _textbox(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(0.5),
             lbl["risk_top"], size=16, bold=True, color=DARK)

    top = df.sort_values("risk_score", ascending=False).head(5)
    y = Inches(2.3)
    for _, row in top.iterrows():
        score = int(round(row["risk_score"]))
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(6.9), y, Inches(0.75), Inches(0.75),
        )
        _fill_shape(oval, RISK_COLORS.get(row["risk_level"], RISK_COLORS["High"]))
        tf = oval.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(score)
        _set_run(run, size=14, bold=True, color=WHITE)

        _textbox(slide, Inches(7.9), y, Inches(4.9), Inches(0.45),
                 row["product_name"], size=14, bold=True)
        sub = lbl["weeks_cover"].format(
            weeks=row["weeks_of_cover"], supplier=row["supplier"],
        )
        _textbox(slide, Inches(7.9), y + Inches(0.4), Inches(4.9), Inches(0.4),
                 sub, size=11, color=MUTED)
        y += Inches(0.87)


def _add_purchase_slide(prs: Presentation, orders: pd.DataFrame,
                        kpis: dict, language: str):
    slide = _blank_slide(prs)
    lbl = LABELS[language]
    _slide_title(slide, lbl["purchase_title"])

    headers = [lbl["col_product"], lbl["col_supplier"], lbl["col_qty"],
               lbl["col_cost"], lbl["col_priority"]]
    rows = orders.copy()
    priority_rank = {"Urgent": 0, "High": 1, "Medium": 2, "Normal": 3}
    rows["_pr"] = rows["priority"].map(priority_rank).fillna(9)
    rows = rows.sort_values(["_pr", "estimated_cost"], ascending=[True, False]).head(8)

    n_rows = max(len(rows), 1) + 1
    table_shape = slide.shapes.add_table(
        n_rows, 5, Inches(0.6), Inches(1.6), Inches(12.1), Inches(4.05),
    )
    table = table_shape.table
    col_widths = [Inches(3.2), Inches(3.5), Inches(1.2), Inches(1.8), Inches(1.4)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                _set_run(run, size=11, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    sorted_orders = rows

    for ri, (_, row) in enumerate(sorted_orders.iterrows(), start=1):
        values = [
            row["product_name"],
            row["supplier"],
            str(int(row["recommended_qty"])),
            f"{row['estimated_cost']:,.0f}",
            row["priority"],
        ]
        for ci, val in enumerate(values):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_run(run, size=10, color=DARK)

    _textbox(
        slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.5),
        lbl["purchase_total"].format(
            value=_fmt_eur(kpis["recommended_purchase_value"]),
        ),
        size=14, bold=True, color=ACCENT,
    )


def _supplier_warning_lines(scorecard: pd.DataFrame, language: str,
                            on_time_threshold: float = 88.0,
                            delay_threshold: float = 3.0) -> list[str]:
    lbl = LABELS[language]
    lines = []
    for _, row in scorecard.iterrows():
        low_otr = row["on_time_rate"] < on_time_threshold
        high_delay = row["avg_delay_days"] > delay_threshold
        delayed = int(row["delayed_deliveries"]) >= 2
        if not (low_otr or high_delay or delayed):
            continue
        params = {
            "supplier": row["supplier"],
            "on_time": row["on_time_rate"],
            "delay": row["avg_delay_days"],
            "delayed": int(row["delayed_deliveries"]),
            "categories": row["categories"],
        }
        if delayed and high_delay:
            template = lbl["supplier_warn"]
        elif delayed:
            template = lbl["supplier_warn_no_delay"]
        else:
            template = lbl["supplier_warn_delay_only"]
        lines.append(template.format(**params))
    return lines[:4]


def _add_supplier_slide(prs: Presentation, scorecard: pd.DataFrame,
                        language: str):
    slide = _blank_slide(prs)
    lbl = LABELS[language]
    _slide_title(slide, lbl["supplier_title"])

    ranked = scorecard.sort_values("on_time_rate", ascending=True).head(8)
    chart_data = CategoryChartData()
    chart_data.categories = [
        s[:22] + "…" if len(s) > 23 else s for s in ranked["supplier"]
    ]
    chart_data.add_series("On-time %", ranked["on_time_rate"].tolist())

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.6), Inches(1.5), Inches(7.0), Inches(5.4),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.value_axis.maximum_scale = 100
    series = chart.plots[0].series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT

    _textbox(slide, Inches(8.0), Inches(1.6), Inches(4.7), Inches(0.5),
             lbl["supplier_attention"], size=16, bold=True, color=DARK)

    warnings = _supplier_warning_lines(scorecard, language)
    y = Inches(2.3)
    for warn in warnings:
        _textbox(slide, Inches(8.0), y, Inches(4.7), Inches(1.0),
                 warn, size=10, color=DARK)
        y += Inches(1.05)


def _add_slow_slide(prs: Presentation, slow: pd.DataFrame, kpis: dict,
                    language: str):
    slide = _blank_slide(prs)
    lbl = LABELS[language]
    _slide_title(slide, lbl["slow_title"])

    if not slow.empty:
        top = slow.sort_values("tied_capital", ascending=False).head(8)
        chart_data = CategoryChartData()
        chart_data.categories = [
            n[:18] + "…" if len(n) > 19 else n for n in top["product_name"]
        ]
        chart_data.add_series("Tied capital", top["tied_capital"].tolist())
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.6), Inches(1.5), Inches(7.6), Inches(5.4),
            chart_data,
        )
        chart = chart_frame.chart
        chart.has_legend = False
        series = chart.plots[0].series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = ACCENT

    _textbox(slide, Inches(8.6), Inches(2.2), Inches(4.1), Inches(1.2),
             _fmt_eur(kpis["slow_mover_capital"]), size=36, bold=True, color=ACCENT)
    _textbox(
        slide, Inches(8.6), Inches(3.3), Inches(4.1), Inches(1.2),
        lbl["slow_capital"].format(count=kpis["slow_mover_count"]),
        size=16, color=MUTED,
    )


def _add_exec_summary_slide(prs: Presentation, summary_text: str, language: str):
    slide = _blank_slide(prs)
    lbl = LABELS[language]
    _slide_title(slide, lbl["exec_summary_title"])

    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.45), Inches(12.1), Inches(4.9),
    )
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for para in summary_text.split("\n\n"):
        para = para.strip()
        if not para:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = para
        _set_run(run, size=14, color=DARK)
        p.space_after = Pt(12)

    _textbox(
        slide, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.45),
        lbl["exec_summary_note"], size=9, color=MUTED,
    )


def _add_actions_slide(prs: Presentation, recs: list[dict], language: str):
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _fill_shape(bg, NAVY)

    lbl = LABELS[language]

    _textbox(slide, Inches(0.6), Inches(0.6), Inches(12.1), Inches(0.9),
             lbl["actions_title"], size=32, bold=True, color=WHITE)

    y = Inches(1.9)
    for i, rec in enumerate(recs[:5], start=1):
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.65), Inches(0.65),
        )
        _fill_shape(oval, ACCENT)
        tf = oval.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i)
        _set_run(run, size=16, bold=True, color=WHITE)

        _textbox(slide, Inches(1.6), y - Inches(0.03), Inches(11.1), Inches(0.5),
                 rec["action"], size=14, bold=True, color=WHITE)
        detail = rec.get("detail", "")
        if detail:
            _textbox(slide, Inches(1.6), y + Inches(0.42), Inches(11.1), Inches(0.55),
                     detail, size=11, color=RGBColor(0xD0, 0xD8, 0xF0))
            y += Inches(0.95)
        else:
            y += Inches(0.72)


def _actions_from_report(report_text: str, language: str) -> list[dict]:
    """Convert numbered management actions from the report into deck bullets."""
    actions = parse_management_actions(report_text, language)
    return [{"action": text, "detail": ""} for text in actions[:5]]


def _deck_recommendations(
    df: pd.DataFrame,
    orders: pd.DataFrame,
    slow: pd.DataFrame,
    scorecard: pd.DataFrame,
    language: str,
    top_n: int = 5,
) -> list[dict]:
    """Localized action bullets for the final deck slide."""
    recs: list[dict] = []

    for _, row in orders[orders["priority"] == "Urgent"].head(3).iterrows():
        days = max(3, int(row["lead_time_days"] // 3))
        if language == "fi":
            action = (
                f"Tilaa {int(row['recommended_qty'])} kpl tuotetta "
                f"{row['product_name']} {days} päivän kuluessa"
            )
            detail = (
                f"Nykyinen varasto {int(row['current_stock'])} kpl, "
                f"toimitusaika {int(row['lead_time_days'])} päivää. "
                f"Arvioitu kustannus {_fmt_eur(row['estimated_cost'])}."
            )
        else:
            action = (
                f"Order {int(row['recommended_qty'])} units of "
                f"{row['product_name']} within {days} days"
            )
            detail = row["reason"] + f" Estimated cost {_fmt_eur(row['estimated_cost'])}."
        recs.append({"action": action, "detail": detail, "urgency": "Urgent"})

    below_safety = df[(df["current_stock"] < df["safety_stock"])
                      & (df["criticality_level"].isin(["high", "critical"]))]
    for _, row in below_safety.head(2).iterrows():
        if language == "fi":
            action = f"Priorisoi {row['product_name']} — alle varmuusvaraston"
            detail = (
                f"Varasto {int(row['current_stock'])} kpl, varmuustaso "
                f"{int(row['safety_stock'])} kpl. Asennuskriittinen tuote; "
                "nopeuta avointa tilausta tai käytä vaihtoehtoista toimittajaa."
            )
        else:
            action = f"Prioritise {row['product_name']} — below safety stock"
            detail = (
                f"Stock {int(row['current_stock'])} vs safety level "
                f"{int(row['safety_stock'])}. Installation-critical item; "
                "expedite the open order or use an alternate supplier."
            )
        recs.append({"action": action, "detail": detail, "urgency": "Urgent"})

    if not slow.empty:
        top_slow = slow.iloc[0]
        if language == "fi":
            action = f"Siirrä {top_slow['product_name']} kampanjasuunnitteluun"
            detail = (
                f"Ylivarastoon sitoutunut {_fmt_eur(top_slow['tied_capital'])}, "
                f"kiertonopeus {top_slow['stock_turnover_rate']:.1f}, "
                f"riitto {top_slow['weeks_of_cover']:.0f} viikkoa."
            )
        else:
            action = f"Move {top_slow['product_name']} to campaign planning"
            detail = (
                f"{_fmt_eur(top_slow['tied_capital'])} tied in excess stock, "
                f"turnover {top_slow['stock_turnover_rate']:.1f}, "
                f"cover {top_slow['weeks_of_cover']:.0f} weeks. "
                f"{top_slow['suggested_action']}"
            )
        recs.append({"action": action, "detail": detail, "urgency": "Medium"})

    if not scorecard.empty and scorecard.iloc[0]["risk_level"] == "High":
        worst = scorecard.iloc[0]
        if language == "fi":
            action = f"Tarkista toimittaja {worst['supplier']}"
            detail = (
                f"Toimitusvarmuus {worst['on_time_rate']:.0f} %, "
                f"keskimääräinen viive {worst['avg_delay_days']:.1f} päivää, "
                f"{int(worst['delayed_deliveries'])} toimitusta myöhässä."
            )
        else:
            action = f"Review supplier {worst['supplier']}"
            detail = (
                f"On-time rate {worst['on_time_rate']:.0f}%, average delay "
                f"{worst['avg_delay_days']:.1f} days, "
                f"{int(worst['delayed_deliveries'])} deliveries currently delayed."
            )
        recs.append({"action": action, "detail": detail, "urgency": "High"})

    urgency_rank = {"Urgent": 0, "High": 1, "Medium": 2}
    recs.sort(key=lambda r: urgency_rank.get(r["urgency"], 3))
    return recs[:top_n]


def export_weekly_deck(
    df: pd.DataFrame,
    orders: pd.DataFrame,
    slow: pd.DataFrame,
    scorecard: pd.DataFrame,
    kpis: dict,
    language: str = "en",
    output_path: str | Path | None = None,
    report_text: str | None = None,
) -> Path:
    """
    Build a weekly deck and save it. Returns the output path.

    When ``report_text`` is provided, inserts an executive-summary slide after
    KPIs and uses the report's management actions on the final slide.
    """
    if language not in LABELS:
        language = "en"

    if output_path is None:
        suffix = "fi" if language == "fi" else "en"
        output_path = Path(f"outputs/reports/weekly_deck_{suffix}.pptx")
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, language)
    _add_kpi_slide(prs, kpis, language)

    exec_summary = ""
    if report_text:
        exec_summary = parse_executive_summary(report_text, language)
    if exec_summary:
        _add_exec_summary_slide(prs, exec_summary, language)

    _add_risk_slide(prs, df, language)
    _add_purchase_slide(prs, orders, kpis, language)
    _add_supplier_slide(prs, scorecard, language)
    _add_slow_slide(prs, slow, kpis, language)

    if report_text:
        deck_recs = _actions_from_report(report_text, language)
    else:
        deck_recs = []
    if not deck_recs:
        deck_recs = _deck_recommendations(df, orders, slow, scorecard, language)
    _add_actions_slide(prs, deck_recs, language)

    prs.save(str(output_path))
    return output_path
